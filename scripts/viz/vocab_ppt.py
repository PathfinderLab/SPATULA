"""Build a presentation PPT explaining the SPATULA vocab build + QC results.

Reads the saved figures + tables from results/eda/prepared_expanded/ and
vocab.csv, lays them out into a slide deck with interpretation text.

Output: results/eda/prepared_expanded/vocab_presentation.pptx
"""
from __future__ import annotations
import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


PREP = Path("/workspace/mm_align/results/cache/prepared_expanded")
EDA = Path("/workspace/mm_align/results/eda/prepared_expanded")
OUT_PPTX = EDA / "vocab_presentation.pptx"

# ── Style helpers ───────────────────────────────────────────────────────────
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x4C, 0x72, 0xB0)
GREEN = RGBColor(0x55, 0xA8, 0x68)
ORANGE = RGBColor(0xDD, 0x88, 0x52)
GRAY = RGBColor(0x55, 0x55, 0x55)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(15), Inches(2))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = NAVY
    # Subtitle
    sb = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(15), Inches(1.5))
    sf = sb.text_frame
    sf.text = subtitle
    sf.paragraphs[0].font.size = Pt(28)
    sf.paragraphs[0].font.color.rgb = GRAY
    return slide


def add_section_slide(prs, title, bullets=None, image_path=None,
                       image_left=8.0, image_top=1.5, image_width=7.5):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Header
    hb = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(15), Inches(0.7))
    htf = hb.text_frame
    htf.text = title
    p = htf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    # Body
    if bullets:
        bb = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(7.5), Inches(7))
        btf = bb.text_frame
        btf.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                para = btf.paragraphs[0]
            else:
                para = btf.add_paragraph()
            txt = b if isinstance(b, str) else b[0]
            sz = 18 if isinstance(b, str) else b[1]
            color = ACCENT if (isinstance(b, str) and b.startswith("✓")) else None
            para.text = txt
            para.font.size = Pt(sz)
            para.font.color.rgb = color or GRAY
            if isinstance(b, str) and b.startswith("✓"):
                para.font.bold = True
                para.font.color.rgb = GREEN
            if isinstance(b, str) and b.startswith("⚠"):
                para.font.color.rgb = ORANGE
                para.font.bold = True
    # Image
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), Inches(image_left), Inches(image_top),
                                  width=Inches(image_width))
    return slide


def add_table_slide(prs, title, dataframe, intro_bullets=None, max_rows=20,
                     col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hb = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(15), Inches(0.7))
    htf = hb.text_frame
    htf.text = title
    p = htf.paragraphs[0]
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = NAVY

    y_top = 1.1
    if intro_bullets:
        ib = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(15), Inches(1.4))
        itf = ib.text_frame; itf.word_wrap = True
        for i, b in enumerate(intro_bullets):
            para = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
            para.text = b
            para.font.size = Pt(16)
            para.font.color.rgb = GRAY
        y_top = 2.4

    df = dataframe.head(max_rows)
    n_rows, n_cols = df.shape
    n_rows += 1   # header

    table_h = min(7.5, 0.30 * n_rows + 0.5)
    if col_widths is None:
        # equal spacing
        total_w = 15.5
        col_widths = [total_w / n_cols] * n_cols
    table_w = sum(col_widths)
    shape = slide.shapes.add_table(n_rows, n_cols,
                                   Inches(0.4), Inches(y_top),
                                   Inches(table_w), Inches(table_h))
    table = shape.table
    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)
    # Header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    # Data
    for i in range(len(df)):
        for j, col in enumerate(df.columns):
            v = df.iloc[i, j]
            if isinstance(v, float):
                s = f"{v:.3g}"
            elif isinstance(v, bool):
                s = "✓" if v else ""
            else:
                s = str(v)
            cell = table.cell(i+1, j)
            cell.text = s
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
            if i % 2 == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF5,0xF5,0xF8)
    return slide


# ── Load data ───────────────────────────────────────────────────────────────
print("Loading vocab + QC outputs…")
vocab = json.loads((PREP / "hvg_vocab.json").read_text())
vocab_df = pd.read_csv(EDA / "vocab_with_tiers.csv")
top20_priority = pd.read_csv(EDA / "vocab_top20_priority.csv")
top20_disp = vocab_df.sort_values("rank").head(20)[
    ["gene","gene_type","sample_prev","spot_prev","nonzero_mean_log1p","norm_dispersion","must_include"]
].round(3)
tier_summary = pd.read_csv(EDA / "tier_summary.csv")

# ── Build deck (16:9) ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Slide 1 — Title
add_title_slide(
    prs,
    "SPATULA Vocab Build & QC",
    f"Principle-driven gene vocabulary for spatial-transcriptomics MSM pretraining\n"
    f"vocab = {len(vocab):,} genes  ·  1,420 shards  ·  1,268 train samples"
)

# Slide 2 — Why we did this
add_section_slide(prs,
    "Why redesign the vocab?",
    bullets=[
        ("Old vocab (n_hvg=4096): seq_len ≈ 8 tokens/spot — model couldn't learn", 20),
        ("Root causes (each cost vocab capacity):", 18),
        ("  · ~52% of vocab was AMBIGUOUS/BAC-clone/pseudogene noise", 16),
        ("  · 10% of HEST (Heart 60, Breast 68) shipped ENSG IDs, never resolved", 16),
        ("  · MT-/RPS/RPL house-keeping silently rescued by GTF guard", 16),
        ("  · No min/max prevalence filter — hapax legomena flooded the rank", 16),
        ("", 12),
        ("Principle-driven approach:", 20),
        ("  1. Strong noise filter + GTF rescue (but never for MT/RPS/RPL)", 16),
        ("  2. ENSG → HGNC resolve at load time (not just spatialcorpus)", 16),
        ("  3. Cross-sample / cross-spot prevalence floors", 16),
        ("  4. Restrict to protein_coding + IG/TR (interpretable)", 16),
        ("  5. n_hvg = CAP, not target — let principles size the vocab", 16),
        ("  6. Curated markers (TP53/MMP/HSP/SFRP/...) force-included", 16),
    ],
)

# Slide 3 — Pipeline
add_section_slide(prs,
    "Pipeline (build_vocab → normalize → zero_removal → token_from_vocab)",
    bullets=[
        ("Stage A — Per-sample load + clean", 18),
        ("  · source dispatch (hest / st1k / spatialcorpus)", 14),
        ("  · symbol normalize: strip prefix/version/suffix, uppercase", 14),
        ("  · ENSG → HGNC via GTF (gencode.v49)", 14),
        ("  · noise pattern drop + protein_coding rescue", 14),
        ("", 10),
        ("Stage B — Global stats (streaming, no full matrix)", 18),
        ("  · per-gene Σx, Σx², Σx_nz, Σ(x_nz)², nz_spots, nz_samples", 14),
        ("  · log1p-normalized values, reservoir sample 1024/gene for median", 14),
        ("", 10),
        ("Stage C — Filter cascade → rank → must_include + (optional) HEG", 18),
        ("", 10),
        ("Stage D — Shard write (raw → normalize_total → log1p → hvg_log)", 18),
        ("", 10),
        ("Stage E — Runtime  (PairedSpotDataset)", 18),
        ("  · gene_norm.mode: nonzero_z (default) or global_median (Geneformer-style)", 14),
        ("  · clip ±8 → Fourier value embedding", 14),
        ("", 10),
        ("Stage F — Zero removal + tokenization (TopHVGGeneEncoder)", 18),
        ("  · drop x=0 positions → variable-length sequence per spot", 14),
        ("  · gene_ids = vocab_token_ids[real positions]", 14),
    ],
)

# Slide 4 — Filter cascade (numbers)
filter_data = pd.DataFrame({
    "Stage": [
        "0. Cleaned + noise + GTF resolve + unmapped ENSG drop",
        "1. min_raw_counts ≥ 3",
        "2. min_sample_prevalence ≥ 0.02",
        "3. min_spot_prevalence ≥ 0.0005",
        "4. gene_type ∈ {protein_coding, IG/TR}",
        "5. Dispersion ranking (n_hvg = null → cap inactive)",
        "Final vocab",
    ],
    "Candidates": [37204, 37204, 23343, 22690, 19183, 19183, 19183],
    "Dropped this step": ["—", 0, "-13,861", "-653", "-3,507", "0 (cap inactive)", "—"],
})
add_table_slide(prs,
    "Filter cascade — 37,204 → 19,183 (principle-driven, no hard cap)",
    filter_data,
    intro_bullets=[
        "Vocab size is a CONSEQUENCE of the principle filters, not a hyper-parameter.",
        "After the cascade, dispersion ranking orders the vocab but doesn't truncate.",
    ],
    max_rows=10,
    col_widths=[8.5, 3.5, 3.5],
)

# Slide 5 — gene_type composition
add_section_slide(prs,
    "Gene-type composition (vocab.csv `gene_type` column)",
    bullets=[
        ("✓ protein_coding: 18,952 (98.8%) ← what we want", 20),
        ("✓ IG/TR rearr. genes: 231 (1.2%) — immune diversity", 18),
        ("✓ Zero lncRNA / pseudogene / unknown — restrict + noise filter worked", 20),
        ("✓ Zero MT-/RPS/RPL — NEVER_RESCUE guard active", 20),
        ("", 12),
        ("Previous iteration (before MT-rescue / ENSG fixes):", 18),
        ("  · unknown 1,609 (38.5%), pseudogene 296 (7.1%)", 14),
        ("  · MT-CO1/2/3, MT-ND4, RPL37A at bottom (high prev × low var)", 14),
    ],
    image_path=EDA / "vocab_qc" / "fig_genetype_pie.png",
)

# Slide 6 — Prevalence
add_section_slide(prs,
    "Prevalence — every vocab gene is now broadly expressed",
    bullets=[
        ("Median sample_prev = 0.925 (vs 0.017 with old vocab — 54× ↑)", 20),
        ("Median spot_prev   = 0.085 (vs 0.0002 — 425× ↑)", 20),
        ("80.5% of vocab in ≥50% of samples", 18),
        ("45.4% of vocab in ≥10% of spots", 18),
        ("", 10),
        ("Why max_*_prevalence is OFF by default:", 16),
        ("  · house-keeping = high prev × LOW variance", 14),
        ("  · dispersion ranking already pushes them down the priority_rank", 14),
        ("  · hard cap would be redundant + brittle", 14),
    ],
    image_path=EDA / "vocab_qc" / "fig_prevalence_scatter.png",
)

# Slide 7 — Marker panel coverage
add_section_slide(prs,
    "Curated marker panels — 152 / 154 captured",
    bullets=[
        ("✓ Tumor suppressors: 16/16 (TP53, RB1, PTEN, BRCA1/2, ...)", 16),
        ("✓ Oncogenes: 16/16 (MYC, KRAS, EGFR, ERBB2, ...)", 16),
        ("✓ MMPs (invasion): 8/8", 16),
        ("✓ Heat shock: 8/8 (HSP90AA1/AB1, HSPA1A/1B, HSPA5, ...)", 16),
        ("✓ Immune T-cell / B-cell / Macrophage: 30/30", 16),
        ("✓ ECM / fibroblast / endothelial / epithelial: 31/31", 16),
        ("✓ Wnt / SFRP antagonists: 16/17 (SFRP3 = alias of FRZB — present)", 16),
        ("", 10),
        ("⚠ Only missing: OCT4 → POU5F1 (alias, present via POU5F1)", 16),
        ("⚠               SFRP3 → FRZB   (alias, present via FRZB)", 16),
    ],
    image_path=EDA / "vocab_qc" / "fig_marker_panel_status.png",
)

# Slide 8 — Top 20 priority genes
add_table_slide(prs,
    "Top 20 by priority_rank — clinical markers lead the vocab",
    top20_priority,
    intro_bullets=[
        "priority = must_include × 10 + in_filter_pool × 5 − rank × 1e-4",
        "First 152 slots are curated markers (must_include=True), then dispersion order.",
    ],
    max_rows=20,
    col_widths=[1.6, 2.3, 1.6, 1.8, 1.6, 2.0, 2.2, 1.8],
)

# Slide 9 — Top 20 by dispersion (organ-specific markers)
add_table_slide(prs,
    "Top 20 by raw dispersion — organ-specific signal at the top",
    top20_disp,
    intro_bullets=[
        "These are NOT in must_include but emerge naturally from variability ranking.",
        "Each entry is a strong organ/tissue-specific marker — exactly what MSM should learn.",
    ],
    max_rows=20,
    col_widths=[1.8, 2.4, 1.8, 1.8, 2.3, 2.4, 2.7],
)

# Slide 10 — Source-level quality
add_section_slide(prs,
    "Quality by data source — Stage 1 token budget per spot",
    bullets=[
        ("seq_len = number of non-zero vocab genes per spot (= sentence length)", 16),
        ("", 10),
        ("✓ HEST: mean 2,105 / median 1,575 / vocab_hit 71%", 18),
        ("✓ spatialcorpus: mean 1,993 / median 1,675 / vocab_hit 77%", 18),
        ("✓ ST1K: mean 2,487 / median 2,021 / vocab_hit 79%", 18),
        ("", 10),
        ("Δ vs previous iteration — HEST mean: 1,774 → 2,105 (+18.7%)", 18),
        ("                          HEST vocab_hit: 0.557 → 0.710 (+27%)", 18),
        ("(spatialcorpus / ST1K unchanged — they already had correct symbols)", 16),
    ],
    image_path=EDA / "process_qc" / "fig_seqlen_by_source.png",
)

# Slide 11 — Organ-level quality
add_section_slide(prs,
    "Quality by organ (HEST) — Heart recovered from process failure",
    bullets=[
        ("⚠ Before fix: Heart seq_len = 31, vocab_hit = 0.014", 18),
        ("✓ After fix:  Heart seq_len = 1,560, vocab_hit = 0.733  (50× / 52× ↑)", 18),
        ("", 10),
        ("Cause: 60/62 Heart samples shipped ENSG IDs (e.g. SPA*, MISC*)", 16),
        ("Fix: _clean_adata_var_names(resolve_ensg=True) sniffs first 32", 16),
        ("     var names, loads GTF lazily, resolves ENSG → HGNC", 16),
        ("", 10),
        ("Same fix recovered 68 Breast samples:", 18),
        ("  · Breast seq_len: 979 → 1,950 (+99%)", 16),
        ("  · Breast vocab_hit: 0.282 → 0.699 (+148%)", 16),
        ("", 10),
        ("Top organs (clean):", 18),
        ("  Lymph node 5,606 · Ovary 4,864 · Cervix 3,656 · Bowel 3,089", 14),
    ],
    image_path=EDA / "process_qc" / "fig_quality_by_organ.png",
    image_top=2.5, image_width=7.5,
)

# Slide 12 — Heart fix table
heart_fix = pd.DataFrame({
    "Organ": ["Heart", "Breast", "Lymph node", "ST1K avg"],
    "n_samples": [62, 125, 5, 672],
    "seq_len_before": [31, 979, 5603, 2487],
    "seq_len_after": [1560, 1950, 5606, 2487],
    "Δ seq_len": ["+50×", "+99%", "+0%", "—"],
    "vocab_hit_before": [0.014, 0.282, 0.862, 0.783],
    "vocab_hit_after": [0.733, 0.699, 0.868, 0.789],
})
add_table_slide(prs,
    "Heart / Breast recovery — before vs after ENSG-resolve fix",
    heart_fix,
    intro_bullets=[
        "60/62 Heart + 68/125 Breast HEST samples were silently broken (ENSG IDs not resolved).",
        "A single-line code fix (resolve_ensg=True) brought them back to par with other organs.",
    ],
    max_rows=10,
    col_widths=[2.3, 1.7, 2.3, 2.3, 2.0, 2.5, 2.5],
)

# Slide 13 — validate_vocab correctness
val_df = pd.read_csv(EDA / "validate_vocab" / "vocab_match_audit.csv")[
    ["sample_id","n_spots_shard","n_spots_source","row_align_ok",
     "n_genes_matched","n_genes_diff","n_genes_zero_filled"]
]
add_table_slide(prs,
    "validate_vocab — shard ↔ source AnnData column-by-column",
    val_df,
    intro_bullets=[
        "For 5 HEST samples we re-load the raw h5ad, run the same prepare pipeline,",
        "and check every vocab column equals the shard's hvg_log column (within 1e-4).",
        "✓ All shards: 0 differences — prepare projection is bit-exact.",
    ],
    max_rows=5,
    col_widths=[2.0, 2.5, 2.5, 2.0, 2.5, 2.0, 2.5],
)

# Slide 14 — Peer-review tiers
add_table_slide(prs,
    "Peer-review tiers — what to keep / what to discuss",
    tier_summary,
    intro_bullets=[
        "review_tier column added to vocab.csv for downstream filtering.",
        "Tier A (must_include + top dispersion) = 2,113 genes — KEEP.",
        "Tier B (rank 2K–8K) = 6,102 genes — KEEP (solid signal).",
        "Tier C (rank 8K–16K) = 7,790 genes — DISCUSS (medium dispersion).",
        "Tier D (rank > 16K) = 3,178 genes — TRIM CANDIDATES if budget is tight.",
    ],
    max_rows=6,
    col_widths=[3.3, 1.5, 2.7, 2.7, 2.0, 2.0, 1.5],
)

# Slide 15 — File locations + peer-review CSV
add_section_slide(prs,
    "Where to find everything (peer-review starting points)",
    bullets=[
        ("📁 Vocab + stats (production)", 20),
        ("  results/cache/prepared_expanded/vocab.csv  ← 15 cols, sorted by priority", 14),
        ("  results/cache/prepared_expanded/hvg_vocab.json  ← used by hvg_log indexing", 14),
        ("  results/cache/prepared_expanded/hvg_vocab_dict.json  ← gene → token_id", 14),
        ("  results/cache/prepared_expanded/gene_stats.npz  ← runtime norm constants", 14),
        ("  results/cache/prepared_expanded/sample_qc.csv  ← per-sample stats", 14),
        ("", 8),
        ("📁 EDA / peer review", 20),
        ("  results/eda/prepared_expanded/vocab_with_tiers.csv ← vocab.csv + review_tier", 14),
        ("  results/eda/prepared_expanded/vocab_top20_priority.csv", 14),
        ("  results/eda/prepared_expanded/tier_summary.csv", 14),
        ("  results/eda/prepared_expanded/tier_genetype_crosstab.csv", 14),
        ("", 8),
        ("📁 Detailed reports", 20),
        ("  results/eda/prepared_expanded/vocab_qc/vocab_qc.md  (+ 7 figures)", 14),
        ("  results/eda/prepared_expanded/process_qc/process_qc.md (+ 4 figures)", 14),
        ("  results/eda/prepared_expanded/validate_vocab/validate_vocab.md", 14),
        ("", 8),
        ("📄 Design + pipeline doc:  docs/design/vocab.md", 20),
    ],
)

# Slide 16 — Next steps
add_section_slide(prs,
    "Next steps — Stage 1 pretraining",
    bullets=[
        ("✓ Vocab is ready — 19,183 genes, all PC/IG/TR, principle-driven", 22),
        ("✓ Shards are bit-exact (validate_vocab pass)", 22),
        ("✓ gene_stats.npz computed (streaming, OOM-safe)", 22),
        ("✓ Tier classification stored in vocab_with_tiers.csv", 22),
        ("", 12),
        ("Launch:  bash scripts/run_stage1.sh", 20),
        ("Config:  configs/stage1/{data,model,train,experiment}.yaml", 18),
        ("", 12),
        ("Open peer-review questions (use vocab_with_tiers.csv):", 22),
        ("  · Keep all Tier C+D, or trim to Tier A+B (8,215 genes)?", 16),
        ("  · Add specific markers we missed? (e.g. lncRNA NEAT1, MALAT1?)", 16),
        ("  · Adjust prevalence thresholds for next iteration?", 16),
    ],
)

prs.save(str(OUT_PPTX))
print(f"\nWrote {OUT_PPTX}  ({OUT_PPTX.stat().st_size/1024:.1f} KB, {len(prs.slides)} slides)")
